from pptx import Presentation

def create_presentation(slides_content):
    prs = Presentation()
    for slide_content in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        print('------',slide_content)
        title.text = slide_content['title']
        body.text = slide_content['body']
    prs.save('presentation.pptx')