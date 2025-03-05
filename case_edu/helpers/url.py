import streamlit as st
import requests
from bs4 import BeautifulSoup
from summa import summarizer
import nltk
from urllib.parse import urlparse
from pptx import Presentation
from io import BytesIO

# Download required NLTK data
nltk.download('punkt')

def fetch_article_content(url):
    """Fetch article content from URL"""
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.text, 'html.parser')
        paragraphs = soup.find_all('p')
        article_text = ' '.join([para.get_text() for para in paragraphs])
        return article_text
    except Exception as e:
        st.error(f"Error fetching article: {str(e)}")
        return None

def create_slides(text, num_slides=5):
    """Create presentation slides from text"""
    summary = summarizer.summarize(text, ratio=0.3)
    sentences = summary.split('. ')
    sentences_per_slide = max(1, len(sentences) // num_slides)
    
    slides = []
    current_slide = {"title": "Introduction", "content": []}
    
    for i, sentence in enumerate(sentences):
        if i % sentences_per_slide == 0 and i != 0:
            slides.append(current_slide)
            current_slide = {"title": f"Key Point {len(slides) + 1}", "content": []}
        current_slide["content"].append(sentence)
    
    slides.append(current_slide)
    return slides

def create_pptx(slides):
    """Create a PowerPoint presentation from slides"""
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Article Presentation"
    subtitle.text = "Generated for Students"
    
    # Content slides
    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_data in slides:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        # Set title
        title_shape = shapes.title
        title_shape.text = slide_data["title"]
        
        # Add content
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        for point in slide_data["content"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    
    # Save to bytes buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def main():
    st.title("Article to Presentation Converter")
    st.write("Turn any article into a student-friendly PowerPoint presentation!")

    # Get URL from user
    url = st.text_input("Enter article URL:", "")
    
    if url:
        # Validate URL
        parsed_url = urlparse(url)
        if not parsed_url.scheme or not parsed_url.netloc:
            st.error("Please enter a valid URL (including http:// or https://)")
            return

        # Process button
        if st.button("Create Presentation"):
            with st.spinner("Generating presentation..."):
                # Fetch article content
                article_text = fetch_article_content(url)
                
                if article_text:
                    # Create slides
                    slides = create_slides(article_text)
                    
                    # Generate PPTX
                    pptx_buffer = create_pptx(slides)
                    
                    # Display slides in app
                    st.success("Presentation generated!")
                    for i, slide in enumerate(slides, 1):
                        st.subheader(slide["title"])
                        content = "\n".join(slide["content"])
                        st.write(content)
                    
                    # Download button
                    st.download_button(
                        label="Download PowerPoint Presentation",
                        data=pptx_buffer,
                        file_name="article_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

if __name__ == "__main__":
    st.set_page_config(page_title="Article Presentation Maker", layout="wide")
    main()